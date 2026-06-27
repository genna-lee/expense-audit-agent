"""
expense_agent/report.py — Monthly Expense Audit Report Generator

CLI usage:
    uv run python -m expense_agent.report
    uv run python -m expense_agent.report --log tests/demo_audit_log.jsonl --month 2026-06 --out reports/

6-page PPTX:
  1. Cover
  2. Overview Statistics
  3. Risk Flag Statistics
  4. Top Suspicious Cases (masked names + case_id)
  5. Security Anomalies (rate-limit + injection events, STRIDE: D · S · T)
  6. Risk Summary & Recommendations (Gemini-generated)

Color scheme:
  - Cover: Navy #1F3864 background, white text
  - Content: White background, dark-grey #2C3E50 text
  - High-risk flags: Red #C0392B
  - Accent: Gold #F0B429
"""

from __future__ import annotations

import argparse
import json
import os
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

# ── python-pptx ──────────────────────────────────────────
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Gemini (google-genai sync SDK) ───────────────────────
from google import genai

from .config import MODEL_NAME

# ── Path constants ────────────────────────────────────────
_DATA_DIR = Path(__file__).parent / "data"
_DEFAULT_LOG = _DATA_DIR / "audit_log.jsonl"
_DEFAULT_OUT = Path(__file__).parent.parent / "reports"

# ── Color palette ─────────────────────────────────────────
C_NAVY   = RGBColor(0x1F, 0x38, 0x64)   # Navy (cover background)
C_GOLD   = RGBColor(0xF0, 0xB4, 0x29)   # Gold (accent)
C_DARK   = RGBColor(0x2C, 0x3E, 0x50)   # Dark grey (body text)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)   # Red (high-risk)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xF5, 0xF6, 0xFA)   # Light grey (card background)
C_BORDER = RGBColor(0xCC, 0xCC, 0xCC)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────
# Utility functions
# ─────────────────────────────────────────────────────────

def mask_name(name: str) -> str:
    """Mask submitter name for privacy: Wang San-Feng -> W○g; John Smith -> J*** S***"""
    if not name or name in ("Unknown", "UNKNOWN"):
        return "○○○"
    # ASCII (usually First Last format)
    if all(c.isascii() for c in name.replace(" ", "")):
        parts = name.split()
        return " ".join(p[0] + "***" if len(p) > 1 else p for p in parts)
    # CJK name
    chars = list(name)
    if len(chars) <= 1:
        return name
    if len(chars) == 2:
        return chars[0] + "○"
    return chars[0] + "○" * (len(chars) - 2) + chars[-1]


def load_log(log_path: Path, month: str | None = None) -> list[dict]:
    """Load audit_log.jsonl, optionally filtered by YYYY-MM."""
    if not log_path.exists():
        return []
    entries = []
    with open(log_path, encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            try:
                entry = json.loads(line)
                if month:
                    date_str = entry.get("date", "") or entry.get("timestamp", "")
                    if not date_str.startswith(month):
                        continue
                entries.append(entry)
            except Exception:
                continue
    return entries


def classify_flags(entries: list[dict]) -> dict[str, int]:
    """Categorize fraud_flags and count by type. Supports both English and Chinese flag text."""
    counts = {
        "Defunct Vendor":    0,
        "Split Purchase":    0,
        "Duplicate Invoice": 0,
        "Over-Budget":       0,
        "Travel Fraud":      0,
        "Injection / PII":   0,
        "Rate Limit":        0,
        "Other":             0,
    }
    for e in entries:
        for flag in e.get("fraud_flags", []):
            f = str(flag)
            if "Defunct" in f or "歇業" in f:
                counts["Defunct Vendor"] += 1
            elif "Split" in f or "拆單" in f:
                counts["Split Purchase"] += 1
            elif "Duplicate Invoice" in f or "重複發票" in f:
                counts["Duplicate Invoice"] += 1
            elif "Over-Limit" in f or "Over-Budget" in f or "業務費" in f or "超上限" in f:
                counts["Over-Budget"] += 1
            elif "Travel" in f or "Inflated" in f or "Hotel" in f or "Misc Over" in f \
                    or "出差" in f or "住宿" in f or "雜費" in f or "浮報" in f:
                counts["Travel Fraud"] += 1
            elif "Injection" in f or "注入" in f or "PII" in f or "REDACTED" in f \
                    or "Security Event" in f:
                counts["Injection / PII"] += 1
            elif "Rate limit" in f or "Account locked" in f:
                counts["Rate Limit"] += 1
            else:
                counts["Other"] += 1
    return counts


def build_stats(entries: list[dict]) -> dict[str, Any]:
    """Aggregate summary statistics."""
    total = len(entries)
    approved = sum(1 for e in entries if e.get("status") == "APPROVED")
    rejected = sum(1 for e in entries if e.get("status") == "REJECTED")
    pending  = total - approved - rejected
    total_amt   = sum(e.get("amount", 0) for e in entries)
    flagged_amt = sum(e.get("amount", 0) for e in entries if e.get("fraud_flags"))
    high_risk       = sum(1 for e in entries if e.get("risk_level") in ("HIGH", "CRITICAL"))
    security_events = sum(1 for e in entries if e.get("risk_level") in ("SECURITY", "CRITICAL"))
    return {
        "total":            total,
        "approved":         approved,
        "rejected":         rejected,
        "pending":          pending,
        "total_amt":        total_amt,
        "flagged_amt":      flagged_amt,
        "high_risk":        high_risk,
        "security_events":  security_events,
        "flag_counts":      classify_flags(entries),
    }


# Category translation map (ZH → EN)
_CAT_EN: dict[str, str] = {
    "設備費":  "Equipment",
    "電腦採購": "Computer Purchase",
    "出差旅費": "Business Travel",
    "辦公設備": "Office Equipment",
    "文具":    "Stationery",
    "雜費":    "Miscellaneous",
    "餐費":    "Meals",
    "交通費":  "Transportation",
    "教育訓練": "Training",
}

def _bilingual_category(cat: str) -> str:
    en = _CAT_EN.get(cat, "")
    if en and en.lower() != cat.lower():
        return f"{cat}\n{en}"
    return cat


def top_suspicious(entries: list[dict], n: int = 20) -> list[dict]:
    """Return flagged cases, deduplicated by case_id, sorted by amount descending."""
    flagged = [e for e in entries if e.get("fraud_flags")]
    # Deduplicate by case_id — prefer REJECTED, then latest timestamp
    seen: dict[str, dict] = {}
    for e in flagged:
        cid = e.get("case_id", "")
        if cid not in seen:
            seen[cid] = e
        else:
            prev = seen[cid]
            if e.get("status") == "REJECTED" and prev.get("status") != "REJECTED":
                seen[cid] = e
            elif (e.get("status") == prev.get("status")
                  and e.get("timestamp", "") > prev.get("timestamp", "")):
                seen[cid] = e
    return sorted(seen.values(), key=lambda e: e.get("amount", 0), reverse=True)[:n]


# ─────────────────────────────────────────────────────────
# PPTX helpers
# ─────────────────────────────────────────────────────────

def _fill_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide, left, top, width, height,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = C_DARK,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
) -> Any:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _style_cell(cell, text: str, font_size: int = 11, bold: bool = False,
                bg: RGBColor | None = None, color: RGBColor = C_DARK,
                align=PP_ALIGN.LEFT) -> None:
    cell.text = text
    tf = cell.text_frame
    tf.word_wrap = True
    # cell.text = "a\nb" creates multiple paragraphs; style ALL of them
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


# ─────────────────────────────────────────────────────────
# 5 slides
# ─────────────────────────────────────────────────────────

def slide_cover(prs: Presentation, month: str, gen_time: str) -> None:
    """Page 1: Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _fill_bg(slide, C_NAVY)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(2.9), SLIDE_W, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_GOLD
    bar.line.fill.background()

    _add_textbox(slide, Inches(1), Inches(0.7), Inches(11), Inches(0.9),
                 "SmartAudit", font_size=48, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(1.55), Inches(11), Inches(0.9),
                 "Expense Audit Report", font_size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(2.45), Inches(11), Inches(0.5),
                 f"Monthly Report  ·  {month}",
                 font_size=20, color=C_GOLD, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.15), Inches(11), Inches(0.5),
                 "AI-Generated by Expense Audit Agent  |  SmartAudit v1.0",
                 font_size=16, color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(3.85), Inches(11.6), Inches(0.55),
                 "Warning: This report is AI-generated and intended for evaluation and internal review only. "
                 "All findings must be verified by an authorized human auditor before any action is taken.",
                 font_size=14, color=RGBColor(0xFF, 0xCC, 0x00), align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(6.85), Inches(11), Inches(0.4),
                 f"Generated: {gen_time}",
                 font_size=12, color=RGBColor(0x95, 0xA5, 0xA6), align=PP_ALIGN.CENTER)


def slide_overview(prs: Presentation, stats: dict, month: str) -> None:
    """Page 2: Overview Statistics"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, C_WHITE)

    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.9))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = C_NAVY; hdr.line.fill.background()
    _add_textbox(slide, Inches(0.3), Inches(0.1), Inches(12), Inches(0.7),
                 f"Overview  |  {month}", font_size=24, bold=True, color=C_WHITE)

    cards = [
        ("Total Claims",    str(stats["total"]),               C_DARK),
        ("Auto-Approved",   str(stats["approved"]),            RGBColor(0x27, 0xAE, 0x60)),
        ("Rejected",        str(stats["rejected"]),            C_RED),
        ("High-Risk Cases", str(stats["high_risk"]),           C_RED),
        ("Total Amount",    f"NT$ {stats['total_amt']:,.0f}",  C_DARK),
        ("Flagged Amount",  f"NT$ {stats['flagged_amt']:,.0f}", C_RED),
    ]
    card_w, card_h = Inches(3.8), Inches(2.2)
    for i, (label, value, val_color) in enumerate(cards):
        col, row = i % 3, i // 3
        left = Inches(0.5) + col * (card_w + Inches(0.3))
        top  = Inches(1.1) + row * (card_h + Inches(0.2))
        rect = slide.shapes.add_shape(1, left, top, card_w, card_h)
        rect.fill.solid(); rect.fill.fore_color.rgb = C_LIGHT
        rect.line.color.rgb = C_BORDER; rect.line.width = Pt(0.5)
        _add_textbox(slide, left + Inches(0.15), top + Inches(0.15),
                     card_w - Inches(0.3), Inches(0.5),
                     label, font_size=16, color=RGBColor(0x7F, 0x8C, 0x8D))
        _add_textbox(slide, left + Inches(0.15), top + Inches(0.65),
                     card_w - Inches(0.3), Inches(1.1),
                     value, font_size=34, bold=True, color=val_color)


def slide_flag_stats(prs: Presentation, stats: dict, month: str) -> None:
    """Page 3: Risk Flag Statistics (horizontal bar chart)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, C_WHITE)

    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.9))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = C_NAVY; hdr.line.fill.background()
    _add_textbox(slide, Inches(0.3), Inches(0.1), Inches(12), Inches(0.7),
                 f"Risk Flag Statistics  |  {month}", font_size=24, bold=True, color=C_WHITE)

    flag_counts = stats["flag_counts"]
    max_val = max(flag_counts.values()) if any(flag_counts.values()) else 1
    bar_max_w = Inches(7.5)
    row_h = Inches(0.9)
    start_top = Inches(1.05)

    high_risk_labels = {"Defunct Vendor", "Split Purchase", "Injection / PII"}
    for i, (label, count) in enumerate(flag_counts.items()):
        top = start_top + i * row_h
        _add_textbox(slide, Inches(0.5), top + Inches(0.18),
                     Inches(3.2), Inches(0.55),
                     label, font_size=16, color=C_DARK, bold=(count > 0))
        if count > 0:
            bar_w = bar_max_w * (count / max_val)
            bar = slide.shapes.add_shape(1, Inches(3.8), top + Inches(0.2),
                                         max(bar_w, Inches(0.3)), Inches(0.45))
            bar_color = C_RED if label in high_risk_labels else C_GOLD
            bar.fill.solid(); bar.fill.fore_color.rgb = bar_color
            bar.line.fill.background()
            _add_textbox(slide, Inches(3.8) + max(bar_w, Inches(0.3)) + Inches(0.12),
                         top + Inches(0.18), Inches(1.2), Inches(0.5),
                         str(count), font_size=20, bold=True, color=bar_color)
        else:
            _add_textbox(slide, Inches(3.8), top + Inches(0.18),
                         Inches(2), Inches(0.5),
                         "0", font_size=16, color=RGBColor(0xBD, 0xC3, 0xC7))


_ROWS_PER_PAGE = 5   # max data rows per P4 slide (excluding header)


def _related_cases_text(entry: dict) -> str:
    """Return related case IDs, excluding the entry's own case_id, one per line."""
    own = entry.get("case_id", "")
    others = [r for r in entry.get("related_case_ids", []) if r and r != own]
    return "\n".join(others) if others else ""


def _flag_summary(entry: dict) -> str:
    """Return full flag summary for table display (Chinese and English preserved)."""
    flags = entry.get("fraud_flags", [])
    return "\n".join(flags) if flags else ""


_COL_DEFS = [
    ("Case ID",               1.5),
    ("Related Cases",         1.5),
    ("Masked Name",           1.0),
    ("Category /\n類別",      1.3),
    ("Amount",                1.2),
    ("Date",                  1.0),
    ("Risk Flag Summary",     4.3),
    ("Status",                1.0),
]


def _render_cases_page(
    prs: Presentation,
    cases: list[dict],
    month: str,
    page_num: int,
    total_pages: int,
) -> None:
    """Render one page of the Top Suspicious Cases table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, C_WHITE)

    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.9))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = C_NAVY; hdr.line.fill.background()
    suffix = f"  ({page_num}/{total_pages})" if total_pages > 1 else ""
    _add_textbox(slide, Inches(0.3), Inches(0.1), Inches(12), Inches(0.7),
                 f"Top Suspicious Cases  |  {month}  (Names Masked){suffix}",
                 font_size=24, bold=True, color=C_WHITE)

    if not cases:
        _add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                     "No suspicious cases this month.", font_size=28,
                     color=RGBColor(0x27, 0xAE, 0x60), align=PP_ALIGN.CENTER)
        return

    n_cols  = len(_COL_DEFS)
    n_rows  = len(cases) + 1
    hdr_h   = Inches(0.55)   # header row
    data_h  = Inches(1.3)    # data row — fits 4 lines of 14pt comfortably
    total_h = hdr_h + data_h * len(cases)
    total_w = Inches(12.8)

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.27), Inches(1.0), total_w, total_h
    ).table

    tbl.rows[0].height = hdr_h
    for ri in range(1, n_rows):
        tbl.rows[ri].height = data_h

    col_widths = [Inches(w) for _, w in _COL_DEFS]
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, (header, _) in enumerate(_COL_DEFS):
        _style_cell(tbl.cell(0, ci), header,
                    font_size=13, bold=True, bg=C_NAVY, color=C_WHITE,
                    align=PP_ALIGN.CENTER)

    for ri, entry in enumerate(cases, start=1):
        is_high = entry.get("risk_level") in ("HIGH", "CRITICAL")
        row_bg  = RGBColor(0xFF, 0xF0, 0xF0) if is_high else C_WHITE
        status  = entry.get("status", "")
        getters = [
            lambda e: e.get("case_id", ""),
            _related_cases_text,
            lambda e: mask_name(e.get("submitter", "")),
            lambda e: _bilingual_category(e.get("category", "")),
            lambda e: f"NT$ {e.get('amount', 0):,.0f}",
            lambda e: e.get("date", ""),
            _flag_summary,
            lambda e: e.get("status", ""),
        ]
        for ci, getter in enumerate(getters):
            val = getter(entry)
            is_status = ci == n_cols - 1
            cell_color = C_RED if (is_status and status == "REJECTED") else C_DARK
            _style_cell(tbl.cell(ri, ci), val,
                        font_size=14,
                        bg=row_bg,
                        color=cell_color,
                        align=PP_ALIGN.CENTER if is_status else PP_ALIGN.LEFT)


def slide_top_cases(prs: Presentation, cases: list[dict], month: str) -> None:
    """Page 4+: Top Suspicious Cases — auto-paginates if > _ROWS_PER_PAGE rows."""
    if not cases:
        _render_cases_page(prs, [], month, 1, 1)
        return
    chunks = [cases[i:i + _ROWS_PER_PAGE]
              for i in range(0, len(cases), _ROWS_PER_PAGE)]
    for idx, chunk in enumerate(chunks, start=1):
        _render_cases_page(prs, chunk, month, idx, len(chunks))


_SEC_COL_DEFS = [
    ("Case ID",    1.8),
    ("Masked Name",1.2),
    ("Date",       1.2),
    ("Type",       1.6),
    ("Detail",     5.8),
    ("Status",     1.2),
]
_SEC_ROWS_PER_PAGE = 6


def _event_type(entry: dict) -> str:
    """Derive security event label from risk_level and flag text."""
    rl = entry.get("risk_level", "")
    flags = [str(f) for f in entry.get("fraud_flags", [])]
    if rl == "CRITICAL":
        return "Injection / PII"
    if any("Account locked" in f for f in flags):
        return "Rate Limit (Hard)"
    if any("Rate limit" in f for f in flags):
        return "Rate Limit (Soft)"
    return rl


def _render_security_page(
    prs: Presentation,
    cases: list[dict],
    month: str,
    page_num: int,
    total_pages: int,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, C_WHITE)

    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.9))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = C_RED; hdr.line.fill.background()
    suffix = f"  ({page_num}/{total_pages})" if total_pages > 1 else ""
    _add_textbox(slide, Inches(0.3), Inches(0.1), Inches(12), Inches(0.7),
                 f"Security Anomalies  |  {month}  |  STRIDE: D · S · T{suffix}",
                 font_size=24, bold=True, color=C_WHITE)

    if not cases:
        _add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                     "No security anomalies detected this month.",
                     font_size=28, color=RGBColor(0x27, 0xAE, 0x60), align=PP_ALIGN.CENTER)
        return

    n_cols  = len(_SEC_COL_DEFS)
    n_rows  = len(cases) + 1
    hdr_h   = Inches(0.55)
    data_h  = Inches(0.95)
    total_h = hdr_h + data_h * len(cases)
    total_w = Inches(12.8)

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.27), Inches(1.0), total_w, total_h,
    ).table

    tbl.rows[0].height = hdr_h
    for ri in range(1, n_rows):
        tbl.rows[ri].height = data_h

    col_widths = [Inches(w) for _, w in _SEC_COL_DEFS]
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, (header, _) in enumerate(_SEC_COL_DEFS):
        _style_cell(tbl.cell(0, ci), header,
                    font_size=13, bold=True, bg=C_RED, color=C_WHITE,
                    align=PP_ALIGN.CENTER)

    C_ORANGE = RGBColor(0xD3, 0x7A, 0x0B)
    for ri, entry in enumerate(cases, start=1):
        is_critical = entry.get("risk_level") == "CRITICAL"
        row_bg  = RGBColor(0xFF, 0xE8, 0xE8) if is_critical else RGBColor(0xFF, 0xF4, 0xE0)
        evt     = _event_type(entry)
        t_color = C_RED if is_critical else C_ORANGE
        flags   = entry.get("fraud_flags", [])
        raw     = flags[0] if flags else ""
        detail  = raw[:120] + "…" if len(raw) > 120 else raw
        status  = entry.get("status", "")
        values = [
            (entry.get("case_id", ""),             C_DARK,                                        PP_ALIGN.LEFT),
            (mask_name(entry.get("submitter", "")), C_DARK,                                        PP_ALIGN.LEFT),
            (entry.get("timestamp", "")[:10],       C_DARK,                                        PP_ALIGN.CENTER),
            (evt,                                   t_color,                                       PP_ALIGN.CENTER),
            (detail,                                C_DARK,                                        PP_ALIGN.LEFT),
            (status,                                C_RED if status == "REJECTED" else C_DARK,     PP_ALIGN.CENTER),
        ]
        for ci, (val, color, align) in enumerate(values):
            _style_cell(tbl.cell(ri, ci), val, font_size=13, bg=row_bg, color=color, align=align)


def slide_security_anomalies(prs: Presentation, entries: list[dict], month: str) -> None:
    """Page 5: Security Anomalies — rate-limit (SECURITY) and injection (CRITICAL) events."""
    sec_entries = [e for e in entries if e.get("risk_level") in ("SECURITY", "CRITICAL")]
    sec_entries.sort(key=lambda e: (
        0 if e.get("risk_level") == "CRITICAL" else 1,
        e.get("timestamp", ""),
    ))
    if not sec_entries:
        _render_security_page(prs, [], month, 1, 1)
        return
    chunks = [sec_entries[i:i + _SEC_ROWS_PER_PAGE]
              for i in range(0, len(sec_entries), _SEC_ROWS_PER_PAGE)]
    for idx, chunk in enumerate(chunks, start=1):
        _render_security_page(prs, chunk, month, idx, len(chunks))


def slide_risk_summary(prs: Presentation, stats: dict, month: str, gemini_text: str) -> None:
    """Page 5: Risk Summary & Recommendations (Gemini-generated)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, C_WHITE)

    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.9))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = C_NAVY; hdr.line.fill.background()
    _add_textbox(slide, Inches(0.3), Inches(0.1), Inches(12), Inches(0.7),
                 f"Risk Summary & Recommendations  |  {month}  ·  AI-Generated",
                 font_size=24, bold=True, color=C_WHITE)

    txt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.0))
    tf = txt_box.text_frame
    tf.word_wrap = True

    for line in gemini_text.split("\n"):
        line = line.strip()
        if not line:
            p = tf.add_paragraph()
            p.space_after = Pt(3)
            continue
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        is_section = line.startswith("【") and "】" in line   # 【摘要】【建議】
        is_zh_rec  = len(line) > 1 and line[1] == "." and line[0].isdigit()  # 1. 中文
        is_en_rec  = line.startswith("   ")                                   # indented English
        if is_section:
            run.font.size  = Pt(17)
            run.font.bold  = True
            run.font.color.rgb = C_NAVY
            p.space_before = Pt(10)
            p.space_after  = Pt(4)
        elif is_zh_rec:
            run.font.size  = Pt(16)
            run.font.bold  = True
            run.font.color.rgb = C_RED
            p.space_before = Pt(6)
            p.space_after  = Pt(1)
        elif is_en_rec:
            run.text = line.strip()
            run.font.size  = Pt(14)
            run.font.bold  = False
            run.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)  # grey
            p.space_after  = Pt(4)
        else:
            run.font.size  = Pt(16)
            run.font.bold  = False
            run.font.color.rgb = C_DARK
            p.space_after  = Pt(4)


# ─────────────────────────────────────────────────────────
# Gemini risk summary
# ─────────────────────────────────────────────────────────

def _load_env() -> None:
    env_file = Path(__file__).parent.parent / ".env"
    if not env_file.exists():
        return
    with open(env_file, encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            k, v = line.split("=", 1)
            k = k.strip()
            v = v.strip().strip('"').strip("'")
            if k and not os.environ.get(k):
                os.environ[k] = v
    if not os.environ.get("GOOGLE_API_KEY") and os.environ.get("GEMINI_API_KEY"):
        os.environ["GOOGLE_API_KEY"] = os.environ["GEMINI_API_KEY"]


def call_gemini_summary(stats: dict, month: str) -> str:
    """Call Gemini to generate a bilingual (Chinese + English) summary and recommendations."""
    _load_env()
    try:
        client = genai.Client()
        prompt = f"""You are a corporate compliance auditor writing a monthly expense audit report for a Taiwan enterprise.
Based on the following audit statistics for {month}, produce a bilingual report in this exact structure:

【本月摘要 Summary】
(2-3 sentences in Traditional Chinese)
(Then the same 2-3 sentences in English)

【改善建議 Recommendations】
1. (繁體中文建議)
   (English translation)
2. (繁體中文建議)
   (English translation)
3. (繁體中文建議)
   (English translation)
4. (繁體中文建議)
   (English translation)
5. (繁體中文建議)
   (English translation)

Output only the above content. No markdown, no extra headers.

Statistics (JSON):
{json.dumps(stats, ensure_ascii=False, indent=2)}
"""
        resp = client.models.generate_content(model=MODEL_NAME, contents=prompt)
        return resp.text.strip()
    except Exception as e:
        print(f"[!] Gemini call failed: {e}")
        total, rejected, high = stats["total"], stats["rejected"], stats["high_risk"]
        return (
            "【本月摘要 Summary】\n"
            f"本月共處理 {total} 筆費用申報，其中 {rejected} 筆遭拒，{high} 筆列為高風險。"
            "主要違規類型包括歇業廠商、拆單採購及出差天數浮報。\n"
            f"This month, {total} expense claims were processed; {rejected} were rejected "
            f"and {high} flagged as high-risk. Key violation types include defunct vendors, "
            "split-purchase evasion, and inflated travel days.\n\n"
            "【改善建議 Recommendations】\n"
            "1. 每月更新歇業廠商名單，付款前強制比對統編。\n"
            "   Refresh the defunct-vendor registry monthly; mandate Tax ID verification before payment.\n"
            "2. 對同一申報人 7 日內多筆採購設定預警，防範拆單規避招標。\n"
            "   Alert on multiple purchases by the same submitter within 7 days to detect split-purchase evasion.\n"
            "3. 出差申請需附行程表與交通憑證，與申報天數交叉比對。\n"
            "   Require itinerary and transport receipts for all travel claims; cross-check against claimed days.\n"
            "4. 在資料輸入層加強提示注入過濾，防止 AI 審核遭操控。\n"
            "   Strengthen prompt-injection filters at data entry to prevent AI auditor manipulation.\n"
            "(Gemini call failed — default bilingual summary displayed)"
        )


# ─────────────────────────────────────────────────────────
# Main flow
# ─────────────────────────────────────────────────────────

def generate_report(
    log_path: Path,
    month: str | None = None,
    out_dir: Path = _DEFAULT_OUT,
) -> Path:
    """Generate PPTX report and return output path."""
    if month is None:
        month = datetime.now(timezone.utc).strftime("%Y-%m")

    print(f"[REPORT] Loading {log_path} (month filter: {month})...")
    entries = load_log(log_path, month)
    if not entries:
        print(f"[REPORT] WARN: No data for {month}; generating empty report.")

    stats    = build_stats(entries)
    top_cas  = top_suspicious(entries)
    gen_time = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")

    print(f"[REPORT] {stats['total']} claims total, {stats['high_risk']} high-risk.")
    print("[REPORT] Calling Gemini for risk summary...")
    gemini_text = call_gemini_summary(stats, month)

    print("[REPORT] Building PPTX...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs, month, gen_time)
    slide_overview(prs, stats, month)
    slide_flag_stats(prs, stats, month)
    slide_top_cases(prs, top_cas, month)
    slide_security_anomalies(prs, entries, month)
    slide_risk_summary(prs, stats, month, gemini_text)

    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"audit_report_{month}.pptx"
    prs.save(str(out_path))
    print(f"[REPORT] Done: {out_path}")
    return out_path


def main() -> None:
    parser = argparse.ArgumentParser(description="Monthly Expense Audit Report Generator")
    parser.add_argument("--log",   type=Path, default=_DEFAULT_LOG)
    parser.add_argument("--month", type=str,  default=None,
                        help="Filter month YYYY-MM (default: current month)")
    parser.add_argument("--out",   type=Path, default=_DEFAULT_OUT,
                        help="PPTX output directory (default: reports/)")
    args = parser.parse_args()
    generate_report(args.log, args.month, args.out)


if __name__ == "__main__":
    main()
